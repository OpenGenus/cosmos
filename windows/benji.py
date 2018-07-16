# coding: utf-8
from __future__ import print_function
import tkinter as tk
from tkinter import ttk
import wx
import regex
import os
import pyautogui
import wikipedia
import time
import webbrowser
import youtube_dl
#import winshell
import json
import requests
import ctypes
import random
import urllib
import datetime
import ssl
from bs4 import BeautifulSoup
import win32com.client as wicl
from urllib.request import urlopen
import speech_recognition as sr
import requests
from pptx import Presentation
from xlsxwriter import Workbook
import subprocess
import sys
import pyttsx3
import getpass
from pytube import YouTube
import matplotlib.pyplot as plt
import numpy as np
import httplib2
import os
from apiclient import discovery
from oauth2client import client
from oauth2client import tools
from oauth2client.file import Storage
import datetime
import face_recognition
import cv2
import tweepy
from tweepy import OAuthHandler
import twitterCredentials
from googletrans import Translator

requests.packages.urllib3.disable_warnings()
try:
		_create_unverified_https_context=ssl._create_unverified_context
except 'AttributeError':
		pass
else:
		ssl._create_default_https_context=_create_unverified_https_context

headers = {'''user-agent':'Chrome/53.0.2785.143'''}
#speak=wicl.Dispatch("SAPI.SpVoice")
speak = pyttsx3.init()

def events(frame,put):
	identity_keywords = ["who are you", "who r u", "what is your name", "who you are"]
	youtube_keywords = ("play ", "stream ", "queue ")
	launch_keywords = ["open ", "launch "]
	search_keywords = ["search "]
	wikipedia_keywords = ["wikipedia ", "wiki "]
	location_keywords = ["locate","spot"]
	check_keywords = ["what","when","was","how","has","had","should","would","can","could","cool","good"] #could or cool or good
	download_music=("download ","download music ")
	search_pc= ("find ","lookfor ")
	graph_generation = ("draw graph for ")
	close_keywords=("close ","over ","stop ","exit ")
	pc_locations = ("desktop", "documents", "downloads")
	
	put = put.lower()
	translator = Translator()
	if put.startswith("spanish"):
		word, space, rest = put.partition(' ')
		translated = translator.translate(rest, src='es', dest='en')
		put = translated.text
	elif put.startswith("french"):
		word, space, rest = put.partition(' ')
		translated = translator.translate(rest, src='fr', dest='en')
		put = translated.text
	link = put.split()

	#Add user for face detection
	if link[0] == "face" or link[0] == "phase":
		name = link[1]
		path = 'C:/dataset' 
		cam = cv2.VideoCapture(0)
		ret, img = cam.read()
		cv2.imwrite(path + "/" + str(name) + ".jpg", img)
		cam.release()
		cv2.destroyAllWindows()
		
	#Get top 10 tweets
	elif link[0] == "get" and link[-1] == "tweets":
		auth = OAuthHandler(twitterCredentials.consumer_key, twitterCredentials.consumer_secret)
		auth.set_access_token(twitterCredentials.access_token, twitterCredentials.access_secret)
		api = tweepy.API(auth)
		if link[-2] == "my":
			for tweet in tweepy.Cursor(api.user_timeline).items(10):
				print("\n", json.dumps(tweet.text))
				print("on ", tweet.created_at)
		elif link[1] == "tweets":
			for status in tweepy.Cursor(api.home_timeline).items(10):
				print("\n", status.text)
				print("By ", status.user.screen_name, " at ", status.user.created_at)
				
	#Get friends from twitter
	elif link[-1] == "twitter" and link[-3] == "follow":
		auth = OAuthHandler(twitterCredentials.consumer_key, twitterCredentials.consumer_secret)
		auth.set_access_token(twitterCredentials.access_token, twitterCredentials.access_secret)
		api = tweepy.API(auth)
		for friend in tweepy.Cursor(api.friends).items():
			print("\nName: ", json.dumps(friend.name), " Username: ", json.dumps(friend.screen_name))		
    #Screenshot    
	elif put.startswith('take screenshot') or put.startswith("screenshot"):
		try:
			pic = pyautogui.screenshot()
			spath = os.path.expanduser('~') + '/Desktop/screenshot.jpg'
			pic.save(spath)
		except:
			print("Unable to take screenshot.")

	#Upcoming events
	elif put.startswith("upcoming events") or put.startswith("coming events") or put.startswith("events"):
		try:
			SCOPES = 'https://www.googleapis.com/auth/calendar.readonly'
			store = file.Storage('credentials.json')
			creds = store.get()
			if not creds or creds.invalid:
				flow = client.flow_from_clientsecrets('client_secret.json', SCOPES)
				creds = tools.run_flow(flow, store)
			service = build('calendar', 'v3', http=creds.authorize(Http()))
				
			now = datetime.datetime.utcnow().isoformat() + 'z' # 'Z' indicates UTC time
			root = tk.Tk()
			root.title("Top 10 Upcoming Events")

			events_result = service.events().list(calendarId='primary', timeMin=now,maxResults=10, singleEvents=True,orderBy='startTime').execute()
			events = events_result.get('items', [])

			if not events:
				w = tk.Label(root, text="No upcoming events found.")
				w.pack()
				
			w = tk.Label(root, text="Event Title")
			w.grid(row=0, column=1)
			w = tk.Label(root, text="Time And Date Of Event")
			w.grid(row=0, column=2)

			i=1
			for event in events:
				start = event['start'].get('dateTime', event['start'].get('date'))
				w = tk.Label(root, text=event['summary'])
				w.grid(row=i, column=1)
				w = tk.Label(root, text=start)
				w.grid(row=i, column=2)
				i=i+1
				
			root.geometry("400x400")
			root.mainloop()
		except:
			print("Unable to take upcoming events")

	#Add note
	elif put.startswith("note") or put.startswith("not") or put.startswith("node"):
		try:
			check = link[1]
			username = os.getlogin()
			filename = "Notes.txt"
			f1 = open(r'''C:\Users\{0}\Desktop\{1}'''.format(username,filename),'a')
			link = '+'.join(link[1:])
			text = link.replace('+',' ')
			text = text[0].capitalize() + text[1:]
			if check in check_keywords:
				text += "?"
			else:
				text += "."
			f1.write(text)
			f1.write("\n")
			f1.close()
			speak.say("Note added successfully!")
			speak.runAndWait()
		except:
			print("Could not add the specified note!")
			
	#adding an event in google calendar
	elif link[0] == "add" and link[1]=="event":
		try:
	            try:
	                import argparse
	                flags = argparse.ArgumentParser(parents=[tools.argparser]).parse_args()
	            except ImportError:
	                flags = None

	            SCOPES = 'https://www.googleapis.com/auth/calendar'
	            CLIENT_SECRET_FILE = 'Client_Secret.json'
	            APPLICATION_NAME = 'GSSOC 	'

	            def get_credentials():
	                home_dir = os.path.expanduser('~')
	                credential_dir = os.path.join(home_dir, '.credentials')
	                if not os.path.exists(credential_dir):
	                    os.makedirs(credential_dir)
	                credential_path = os.path.join(credential_dir,'calendar-python-quickstart.json')
	                store = Storage(credential_path)
	                credentials = store.get()
	                if not credentials or credentials.invalid:
	                    flow = client.flow_from_clientsecrets(CLIENT_SECRET_FILE, SCOPES)
	                    flow.user_agent = APPLICATION_NAME
	                    if flags:
	                        credentials = tools.run_flow(flow, store, flags)
	                    else:
	                        credentials = tools.run(flow, store)
	                    print('Storing credentials to ' + credential_path)
	                return credentials

	            def main():
	                credentials = get_credentials()
	                http = credentials.authorize(httplib2.Http())
	                service = discovery.build('calendar', 'v3', http=http)
	                summary = link[2]
	                d = link[-3]
	                e = link[-1]
	                date = d+"T00:00:00-07:00"
	                end = e+"T00:00:00-07:00"
	                event = {
				'summary': summary,
				'start': {
					'dateTime': date,
			    },
			  'end': {
			    'dateTime': end,
			    },
			  'reminders': {
			    'useDefault': False,
			    'overrides': [
			      {'method': 'email', 'minutes': 24 * 60},
			      {'method': 'popup', 'minutes': 15},
			    ],
			  },
			}

	                event = service.events().insert(calendarId='primary', body=event).execute()
                    #print('Event created: %s' % (event.get('htmlLink')))
					#webbrowser.open('https://calendar.google.com/calendar/r')

	            if __name__ == '__main__':
	                main()

		except Exception as e:
	            print(e)
	#Open a existing folder
	elif put.startswith(search_pc):
		try:
			if any(word in put for word in pc_locations):
				username = getpass.getuser()
				location = link[-1]
				file_name = link[1]
				path = r"C:\Users\%s\%s\%s" %( username, location, file_name)
				os.system("start "+path)
			elif link[-1] == "drive" and link[-3] == "in":
				drive = link[-2]
				file_name1 = link[1]
				if link[2] == link[-3]:
					file_name2 = ''
				else:
					file_name2 = link[2]
				path = r"%s:\%s %s " %(drive, file_name1, file_name2)
				os.system("start " +path)
			elif link[-1] == "drive":
				drive = link[-2]
				path = r"%s:\ " %(drive)
				os.system("start "+path)
		except Exception as e:
			print(e)
	#Screen Recorder
	elif link[0] == "recorder":
		try:
			if len(link) < 2:
				video = '"UScreenCapture"'
				audio = '"Microphone (Realtek High Definition Audio)"'
			elif len(link) < 3:
				video = link[1]
				video = video.replace('_',' ')
				video = '"' + video + '"'
				audio = '"Microphone (Realtek High Definition Audio)"'
			else:
				video = link[1]
				video = video.replace('_',' ')
				video = '"' + video + '"'
				audio = link[2]
				audio = audio.replace('_',' ')
				audio = '"' + audio + '"'
			username = os.getlogin()
			speak.say("Recording started!")
			speak.runAndWait()
			os.chdir(r'''C:\Users\{}\Desktop'''.format(username))
			subprocess.call(r'''ffmpeg -rtbufsize 1500M -f dshow -i video={0}:audio={1} -vcodec mpeg4 -vtag xvid -qscale:v 0 -crf 0 -acodec libmp3lame -ab 320k -ac 1 -ar 44100 video.avi'''.format(video,audio),shell=True)    #video = UScreenCapture , audio = Microphone (Realtek High Definition Audio)
		except:
			print("Unable to start requested service!")
	#Voice Recorder
	elif link[0] == "audio" and link[1] == "recorder":
		try:
			if len(link) < 3:
				audio = '"Microphone (Realtek High Definition Audio)"'
			else:
				audio = link[2]
				audio = audio.replace('_',' ')
				audio = '"' + audio + '"'
			username = os.getlogin()
			speak.say("Recording started!")
			speak.runAndWait()
			os.chdir(r'''C:\Users\{}\Desktop'''.format(username))
			subprocess.call(r'''ffmpeg -rtbufsize 1500M -f dshow -i audio={0} -acodec libmp3lame -ab 320k -ac 1 -ar 44100 audio.mp3'''.format(audio),shell=True)
		except:
			print("Unable to start requested service!")
	#Video Recorder
	elif link[0] == "video" and link[1] == "recorder":
		try:
			if len(link) < 3:
				video = '"UScreenCapture"'
			else:
				video = link[2]
				video = video.replace('_',' ')
				video = '"' + video + '"'
			username = os.getlogin()
			speak.say("Recording started!")
			speak.runAndWait()
			os.chdir(r'''C:\Users\{}\Desktop'''.format(username))
			subprocess.call(r'''ffmpeg -rtbufsize 1500M -f dshow -i video={0} -vcodec mpeg4 -vtag xvid -qscale:v 0 -crf 0 video.avi'''.format(video),shell=True)
		except:
			print("Unable to start requested service!")
	#Merge audio and video
	elif link[0] == "merge":
		try:
			username = os.getlogin()
			os.chdir(r'''C:\Users\{}\Desktop'''.format(username))
			video = link[1]
			audio = link[2]
			output = link[3]
			subprocess.call(r'''ffmpeg -i {} -i {} -c:v copy -c:a copy {}'''.format(video,audio,output),shell=True)
		except:
			print("Unable to process requested service!")
	#Convert video
	elif link[0] == "convert":
		try:
			username = os.getlogin()
			os.chdir(r'''C:\Users\{}\Desktop'''.format(username))
			if link[1] == "na":
				form_in = link[2]
				video1 = link[3]
				form_out = link[4]
				video2 = link[5]
				if (form_in == "avi" or form_in == "webm" or form_in == "mp4" or form_in == "mkv") and (form_out == "mp4" or form_out == "mkv"):
					subprocess.call(r'''ffmpeg -i {} -c:v libx264 -an {}'''.format(video1,video2), shell = True)
				elif (form_in == "avi" or form_in == "mp4" or form_in == "mkv") and form_out == "webm":
					subprocess.call(r'''ffmpeg -i {} -c:v libvpx-vp9 -b:v 2M -an {}'''.format(video1,video2),shell=True)
			else:
				form_in = link[1]
				video1 = link[2]
				form_out = link[3]
				video2 = link[4]
				if (form_in == "avi" or form_in == "webm" or form_in == "mp4" or form_in == "mkv") and (form_out == "mp4" or form_out == "mkv"):
					subprocess.call(r'''ffmpeg -i {} -c:v libx264 -acodec aac {}'''.format(video1,video2), shell = True)
				elif (form_in == "avi" or form_in == "mp4" or form_in == "mkv") and form_out == "webm":
					subprocess.call(r'''ffmpeg -i {} -c:v libvpx-vp9 -b:v 2M -cpu-used -5 -deadline realtime -c:a libvorbis {}'''.format(video1,video2), shell = True)
				elif (form_in == "mp4" or form_in == "mkv" or form_in == "webm") and form_out == "avi":
					subprocess.call(r'''ffmpeg -i {} -c:v mpeg4 -vtag xvid -qscale:v 0 -acodec libmp3lame {}'''.format(video1,video2), shell = True)
				elif (form_in == "avi" or form_in == "webm" or form_in == "mp4" or form_in == "mkv" or form_in == "mp3" or form_in == "m4a") and (form_out == "m4a" or form_out == "mp3"):
					subprocess.call(r'''ffmpeg -i {} {}'''.format(video1,video2), shell = True)
		except:
			print("Unable to process requested service!")

	#Closing Benji
	elif put.startswith(close_keywords):
		os._exit(0)


	#Images to video
	elif put.startswith("images to video "):
		try:
			framerate = link[3]
			username = os.getlogin()
			os.chdir(r'''C:\Users\{}\Desktop\Images'''.format(username))
			subprocess.call(r'''ffmpeg -framerate 1/{} -i img%03d.jpg -vcodec mpeg4 -vtag xvid -qscale:v 0 -crf 0 output.avi'''.format(framerate),shell=True)
			speak.say("Video created!")
			speak.runAndWait()
		except:
			print("Unable to create video file!")

	#Open Files
	elif put.startswith(search_pc):
		try:
			name=link[1]
			rex=regex.compile(name)
			filepath=link[2]
			realpath=filepath
			for root,dirs,files in os.walk(os.path.normpath(filepath)):
				for f in files:
					result = rex.search(f)
					if result:
						realpath=os.path.join(root, f)
						print (realpath+"\n")
			os.startfile(realpath)
		except:
			print("Error")

	#Plotting of graph 
	elif put.startswith(graph_generation):
		try:
			formula = link[3]
			lower_limit = int(link[5])
			upper_limit = int(link[7])
			x = np.array(range(lower_limit,upper_limit))
			y = eval(formula)
			speak.say("Plotting The Graph")
			speak.runAndWait()
			plt.plot(x, y)
			plt.show()
		except:
			print("Error")
			speak.say("Sorry Graph can not be Plotted")
			speak.runAndWait()

#    elif put.startswith(search_pc):
#        process=subprocess.Popen("dir /b/s "+link[1],shell=True,stdout=subprocess.PIPE)
#        while True:
#            output = process.stdout.readline()
#            if output == '' and process.poll() is not None:
#                break
#            if output:
#                print (output.strip()+"\n")
#                outp=output
#        try:
#            os.startfile(outp)
#        except:
#            speak.say("Sorry,couldn't open")

	#Play song on youtube
	if put.startswith(youtube_keywords):
		try:
			link = '+'.join(link[1:])
#                   print(link)
			say = link.replace('+', ' ')
			url = 'https://www.youtube.com/results?search_query='+link
#                 webbrowser.open('https://www.youtube.com'+link)
			fhand=urllib.request.urlopen(url).read()
			soup = BeautifulSoup(fhand, "html.parser")
			songs = soup.findAll('div', {'class': 'yt-lockup-video'})
			hit = songs[0].find('a')['href']
#                   print(hit)
			speak.say("playing "+say)
			speak.runAndWait()
			webbrowser.open('https://www.youtube.com'+hit)
		except:
			print('Sorry Ethan. Looks like its not working!')
	#Download video
	if put.startswith("download video "):
		try:
			link = '+'.join(link[2:])
			say = link.replace('+', ' ')
			url = 'https://www.youtube.com/results?search_query='+link
			fhand=urllib.request.urlopen(url).read()
			soup = BeautifulSoup(fhand, "html.parser")
			songs = soup.findAll('div', {'class': 'yt-lockup-video'})
			hit = songs[0].find('a')['href']
			speak.say("downloading video "+say)
			speak.runAndWait()
			username = os.getlogin()
			os.chdir(r'''C:\Users\{}\Desktop'''.format(username))
			YouTube('https://www.youtube.com' + hit).streams.first().download()
			speak.say("download complete!")
			speak.runAndWait()
		except:
			print('Sorry Ethan. Looks like its not working!')
	#Download music
	elif put.startswith(download_music):
		try:
			link = '+'.join(link[1:])
#                   print(link)
			say = link.replace('+', ' ')
			url = 'https://www.youtube.com/results?search_query='+link
#                 webbrowser.open('https://www.youtube.com'+link)
			fhand=urllib.request.urlopen(url).read()
			soup = BeautifulSoup(fhand, "html.parser")
			songs = soup.findAll('div', {'class': 'yt-lockup-video'})
			hit = songs[0].find('a')['href']
#                   print(hit)
			speak.say("downloading "+say)
			speak.runAndWait()
			ydl_opts = {
						'format': 'bestaudio/best',
						'postprocessors': [{
											'key': 'FFmpegExtractAudio',
											'preferredcodec': 'mp3',
											'preferredquality': '192',
											}],
											'quiet': True,
											'restrictfilenames': True,
											'outtmpl': 'C:\\Users\\'+os.environ['USERNAME']+'\\Desktop\\%(title)s.%(ext)s'
											}

			ydl = youtube_dl.YoutubeDL(ydl_opts)
			ydl.download(['https://www.youtube.com'+hit])
			speak.say("download completed Check your desktop for the song")
			speak.runAndWait()
		except:
			print("Unable to download requested music!")
	#Location
	elif any(word in put for word in location_keywords):
		try:
			link='+'.join(link[1:])
			say=link.replace('+',' ')
			speak.say("locating "+ say)
			speak.runAndWait()
			webbrowser.open('https://www.google.nl/maps/place/'+link)
		except:
			print('The place seems to be sequestered.')
	#Who are you?
	elif any(word in put for word in identity_keywords):
		try:
			speak.say("I am BENJI, a digital assistant declassified for civilian use. Previously I was used by the Impossible Missions Force")
			speak.runAndWait()
		except:
			print('Error. Try reading the ReadMe to know about me!')
	#Open a webpage
	elif any(word in put for word in launch_keywords):
		try:
			link = '+'.join(link[1:])
			speak.say("opening "+link)
			speak.runAndWait()
			webbrowser.open('http://www.'+ link)
		except:
			print('Sorry Ethan,unable to access it. Cannot hack either-IMF protocol!')
	#Google search
	elif any(word in put for word in search_keywords):
		try:
			link='+'.join(link[1:])
			say=link.replace('+',' ')
			speak.say("searching google for "+say)
			speak.runAndWait()
			webbrowser.open('https://www.google.com/search?q='+link)
		except:
			print('Nope, this is not working.')
	#Google Images
	elif put.startswith("images of "):
		try:
			link='+'.join(link[2:])
			say=link.replace('+',' ')
			speak.say("searching images of " + say)
			speak.runAndWait()
			webbrowser.open('https://www.google.co.in/search?q=' + link + '&source=lnms&tbm=isch')
		except:
			print('Could not search for images!')
	#Gmail
	elif put.startswith("gmail"):
		try:
			speak.say("Opening Gmail!")
			speak.runAndWait()
			webbrowser.open('https://www.google.com/gmail')
		except:
			print("Could not open Gmail!")
	#Google Cloud Print
	elif put.startswith("google cloud print"):
		try:
			speak.say("Opening google cloud print!")
			speak.runAndWait()
			webbrowser.open('https://www.google.com/cloudprint')
		except:
			print("Could not open Google Cloud Print!")
	#Google Others
	elif put.startswith("google "):
		try:
			say = link[1]
			speak.say("Opening google " + say)
			speak.runAndWait()
			webbrowser.open('https://'+ say +'.google.com')
		except:
			print("Could not open Google " + say.capitalize() + "!")
	#Blogger
	elif put.startswith("blogger"):
		try:
			speak.say("Opening blogger!")
			speak.runAndWait()
			webbrowser.open('https://www.blogger.com')
		except:
			print("Could not open Blogger!")
	#Wikipedia
	elif any(word in put for word in wikipedia_keywords):
		try:
			link = '+'.join(link[1:])
			say = link.replace('+', ' ')
			wikisearch = wikipedia.page(say)
			speak.say("Opening wikipedia page for" + say)
			speak.runAndWait()
			webbrowser.open(wikisearch.url)
		except:
			print('Wikipedia could not either find the article or your Third-world connection is unstable')
	#Podcast
	elif put.startswith("podcast"):
		try:
			speak.say("Opening podcast!")
			speak.runAndWait()
			webbrowser.open('https://castbox.fm/home')
		except:
			print("Could not open podcast!")
	#Lock the device
	elif put.startswith('secure ') or put.startswith('lock '):
		try:
			speak.say("locking the device")
			speak.runAndWait()
			ctypes.windll.user32.LockWorkStation()
		except :
			print('Cannot lock device')
	#News of various press agencies
	elif put.startswith('news '):
		try:
			say = '+'.join(link[1:])
			say = say.replace('+','-')
			if link[1] == "al" and link[2] == "jazeera":
				say += "-english"
			elif link[1] == "bbc":
				say += "-news"
			elif link[1] == "espn" and link[2] == "cric":
				say += "-info"
			url = ('https://newsapi.org/v1/articles?source=' + say + '&sortBy=latest&apiKey=571863193daf421082a8666fe4b666f3')
			newsresponce = requests.get(url)
			newsjson = newsresponce.json()
			speak.say('Our agents from ' + say + ' report this')
			speak.runAndWait()
			print('  ====='+ say.upper() +'===== \n')
			i = 1
			for item in newsjson['articles']:
				print(str(i) + '. ' + item['title'] + '\n')
				print(item['description'] + '\n')
				i += 1
		except:
			print('Unable to retrieve data!')
	#shutdown after specific time
	elif put.startswith('shutdown after '):
		try:
			if link[2].isdigit() and link[4].isdigit():
				if link[2] == "zero":
					link[2] = "0"
				if link[4] == "zero":
					link[4] = "0"
				hours = int(link[2])
				minutes = int(link[4])
				time_seconds = 60 * minutes
				time_seconds = time_seconds + hours * 3600
				subprocess.call("shutdown /s /t {0}".format(str(time_seconds)), shell = True)
				speak.say("Shutdown initialized!")
				speak.runAndWait()
		except:
			print("Please shutdown manually!")
	#shutdown now
	elif put.startswith("shutdown now"):
		try:
			subprocess.call("shutdown /s /t 0", shell = True)
		except:
			print("Please shutdown manually!")
	#abort shutdown
	elif put.startswith("cancel shutdown"):
		try:
			subprocess.call("shutdown /a", shell = True)
			speak.say("Shutdown cancelled!")
			speak.runAndWait()
		except:
			print("Unable do cancel shutdown!")
	#restart
	elif put.startswith("restart now"):
		try:
			subprocess.call("shutdown /r /t 0", shell = True)
		except:
			print("Unable do restart device!")
	#Folder
	elif put.startswith('create ') and link[-1] == "folder":
		try:
			username = os.getlogin()
			filename = '+'.join(link[1:-1])
			filename = filename.replace('+','_').capitalize()
			path = r'''C:\Users\{0}\Desktop\{1}'''.format(username,filename)
			os.mkdir(path)
			speak.say("Folder created!")
			speak.runAndWait()
		except:
			print("Couldn't create specified folder!")
	#create file
	elif put.startswith('create ') and link[-1] == "document":
		try:
			username = os.getlogin()
			filename = '+'.join(link[1:-2])
			filename = filename.replace('+','_').capitalize()
			if link[-2] == "text":
				filename += ".txt"
				f1 = open(r'''C:\Users\{0}\Desktop\{1}'''.format(username,filename),'a')
				f1.close()
			elif link[-2] == "word" or link[-2] == "world":
				filename += ".docx"
				f1 = open(r'''C:\Users\{0}\Desktop\{1}'''.format(username,filename),'a')
				f1.close()
			elif link[-2] == "powerpoint" or link[-2] =="presentation":
				filename += ".pptx"
				prs = Presentation()
				title_slide_layout = prs.slide_layouts[0]
				slide = prs.slides.add_slide(title_slide_layout)
				os.chdir(r'''C:\Users\{0}\Desktop'''.format(username))
				prs.save(filename)
			elif link[-2] == "excel" or link[-2] == "Excel":
				filename += ".xlsx"
				wb = Workbook(filename)
				ws = wb.add_worksheet()
				os.chdir(r'''C:\Users\{0}\Desktop'''.format(username))
				wb.close()
			elif link[-2] == "visio" or link[-2] == "vizio":
				filename += ".vsdx"
				f1 = open(r'''C:\Users\{0}\Desktop\{1}'''.format(username,filename),'a')
				f1.close()
			elif link[-2] == "rich" or link[-2] == "reach":
				filename += ".rtf"
				f1 = open(r'''C:\Users\{0}\Desktop\{1}'''.format(username,filename),'a')
				f1.close()
			speak.say("Created" + filename)
			speak.runAndWait()
		except:
			print("Unable to create a file.")
	#Calculator
	elif put.startswith('calculator'):
		try:
			subprocess.call('calc',shell=True)
		except:
			print("Unable to open calculator!")
	#Exit/Quit
	elif put.startswith('exit') or put.startswith('quit'):
		sys.exit()

#A stdout class to redirect output to tkinter window
class StdRedirector(object):

	def __init__(self, text_window):
		self.text_window = text_window

	def write(self, output):
		self.text_window.insert(tk.END, output)
# Creating the graphical user interface
class MyFrame(tk.Frame):
	def __init__(self,*args,**kwargs):

		self.textBox = tk.Text(root,
			height=1,width=30,
			font=("Times", 16),
			bg="#666", fg="#0f0",
			spacing1=6, spacing3=6,
			insertbackground="#0f0"
			)
		self.textBox.insert("1.0", "$>")
		self.textBox.grid(row=1,column=1, padx=10, pady=10)
		root.bind('<Return>', self.OnEnter)
		#root.bind('<Destroy>', self.onClose)
		self.textBox.focus_set()
		speak.say('''Hi Agent! BENJI at your service''')
		speak.runAndWait()

		self.photo1 = tk.PhotoImage(file="mic_icon.png")

		self.btn = ttk.Button(root,command=self.OnClicked,
		image=self.photo1, style="C.TButton")
		self.btn.grid(row=1,column=2, padx=10, pady=20)

		'''
		self.output_window = tk.Toplevel()
		output_text_window = tk.Text(self.output_window)
		self.stddirec = StdRedirector(output_text_window)
		sys.stdout = self.stddirec
		output_text_window.pack()
		self.output_window.withdraw()
		'''

	def OnEnter(self,event):
			put=self.textBox.get("1.2","end-1c")
			self.displayText(put)
			self.textBox.insert('1.2',put)
			self.textBox.delete('1.2',tk.END)
			events(self, put)
			if put=='':
			   self.displayText('Reenter')

	def OnClicked(self):
		r = sr.Recognizer()
		with sr.Microphone() as source:
			speak.say('Hey I am Listening ')
			speak.runAndWait()
			audio = r.listen(source)
		try:
			put=r.recognize_google(audio)
			self.displayText(put)
			self.textBox.insert('1.2',put)
			self.textBox.delete('1.2',tk.END)
			events(self,put)
		except sr.UnknownValueError:
			self.displayText("Could not understand audio")
		except sr.RequestError as e:
			self.displayText("Could not request results; {0}".format(e))

	def displayText(self, text):
		try :
			if not self.output_window.winfo_viewable() :
				self.output_window.update()
				self.output_window.deiconify()
		except :
			self.createOutputWindow()
		print(text)

	def createOutputWindow(self):
		self.output_window = tk.Toplevel()
		output_text_window = tk.Text(self.output_window)
		self.stddirec = StdRedirector(output_text_window)
		sys.stdout = self.stddirec
		output_text_window.pack()

	#Trigger the GUI. Light the fuse!
if __name__=="__main__":

	#Face detection
	path = 'C:/'
	lisdir = os.listdir(path)
	flag = 0
	for lis in lisdir:
		# if users face is in dataset, then the following code will run for authentication
		if lis == 'dataset':
			face_recognition.main()
			flag = 1

	if flag != 1:		
		os.mkdir('C:/dataset')
		name = "admin"
		path = 'C:/dataset'
		cam = cv2.VideoCapture(0)
		ret, img = cam.read()
		cv2.imwrite(path + "/" + str(name) + ".jpg", img)
		cam.release()
		cv2.destroyAllWindows()

	#GUI    
	root = tk.Tk()
	view = MyFrame(root)
	style = ttk.Style()
	style.configure('C.TButton',
		background='#555',
		highlightthickness='0'
	)
	style.map("C.TButton",
		background=[('pressed', '!disabled', '#333'), ('active', '#666')]
	)
	# root.geometry('{}x{}'.format(400, 100))
	# view.pack(side="top",fill="both",expand=False)
	root.iconphoto(True, tk.PhotoImage(file=os.path.join(sys.path[0],'benji_final.gif')))
	root.title('B.E.N.J.I.')
	root.configure(background="#444")
	root.resizable(0,0)
	root.mainloop()
